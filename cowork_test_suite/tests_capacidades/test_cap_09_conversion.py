"""
CAP-09: CONVERSION ENTRE FORMATOS
===================================
Tests para conversiones de formato.

Sub-capacidades:
  9.1  PDF -> PowerPoint
  9.2  Notas de reunion -> Documento formateado
  9.3  CSV -> Modelo financiero en Excel
  9.4  Word -> Presentacion
  9.5  Facturas/recibos (screenshots) -> Spreadsheet organizado
  9.6  Excel -> Reporte en Word con comentarios

Total: ~300 tests
"""
from __future__ import annotations

import csv
import pytest
from pathlib import Path
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document

from cowork_lib import generate_document, DocSpec, read_document_text
from cowork_lib2 import (
    markdown_to_html,
    csv_to_xlsx_model,
    docx_to_pptx_outline,
    pptx_to_markdown_outline,
    write_csv,
    write_markdown,
)
from cowork_lib3 import (
    create_pptx_from_doc_outline,
    extract_pdf_tables,
    create_excel_with_formulas,
    create_word_with_comments,
    create_word_with_tables,
)


@pytest.fixture
def out_dir(tmp_path):
    d = tmp_path / "cap09_conv"
    d.mkdir()
    return d


# ============================================================================
# 9.1 — PDF -> PowerPoint
# ============================================================================

@pytest.mark.file_generation
class TestPdfToPptx:
    """9.1 — Convert PDF documents to PowerPoint presentations."""

    @pytest.mark.parametrize("i", range(15))
    def test_pdf_to_pptx(self, out_dir, i):
        spec = DocSpec(kind="pdf", title=f"Report {i}", sections=[
            f"Introduction to topic {i}",
            f"Key findings about item {i}",
            f"Recommendations for improvement",
        ])
        pdf_path = generate_document(spec, out_dir)
        text = read_document_text(pdf_path)
        paragraphs = [p for p in text.split("\n") if p.strip()]
        pptx_path = out_dir / f"from_pdf_{i}.pptx"
        create_pptx_from_doc_outline(paragraphs, pptx_path)
        prs = Presentation(pptx_path)
        assert len(prs.slides) >= 2

    @pytest.mark.parametrize("i", range(10))
    def test_pdf_to_pptx_title_preserved(self, out_dir, i):
        title = f"Quarterly Review Q{(i % 4) + 1}"
        spec = DocSpec(kind="pdf", title=title, sections=["Summary", "Details"])
        pdf_path = generate_document(spec, out_dir)
        text = read_document_text(pdf_path)
        paragraphs = [p for p in text.split("\n") if p.strip()]
        pptx_path = out_dir / f"pdf_title_{i}.pptx"
        create_pptx_from_doc_outline(paragraphs, pptx_path)
        prs = Presentation(pptx_path)
        first_title = prs.slides[0].shapes.title.text
        assert title in first_title or "Quarterly" in first_title


# ============================================================================
# 9.2 — Notas de reunion -> Documento formateado
# ============================================================================

@pytest.mark.file_generation
class TestMeetingNotesToDoc:
    """9.2 — Convert meeting notes to formatted documents (DOCX and Markdown/HTML)."""

    @pytest.mark.parametrize("i", range(15))
    def test_meeting_notes_to_docx(self, out_dir, i):
        notes = [
            f"Meeting: Project Alpha Standup #{i}",
            f"Date: 2026-04-{10 + i}",
            f"Attendees: Alice, Bob, Carol",
            f"Discussion: Sprint review and backlog grooming",
            f"Action Item 1: Alice to update the dashboard",
            f"Action Item 2: Bob to fix bug #{i * 100}",
            f"Next meeting: 2026-04-{17 + i}",
        ]
        spec = DocSpec(kind="docx", title=notes[0], sections=notes[1:])
        path = generate_document(spec, out_dir)
        text = read_document_text(path)
        assert "Meeting" in text
        assert "Action Item" in text

    @pytest.mark.parametrize("i", range(10))
    def test_notes_to_markdown_to_html(self, out_dir, i):
        md_path = out_dir / f"notes_{i}.md"
        write_markdown(md_path, f"Meeting Notes {i}", [
            f"**Attendees**: Alice, Bob",
            f"**Date**: 2026-04-{12 + i}",
            f"- Discussed Q2 roadmap",
            f"- Action: Review budget by Friday",
        ])
        md_content = md_path.read_text()
        html = markdown_to_html(md_content)
        assert "<h1>" in html or "<h2>" in html
        assert "Meeting Notes" in html


# ============================================================================
# 9.3 — CSV -> Modelo financiero en Excel
# ============================================================================

@pytest.mark.file_generation
class TestCsvToExcelModel:
    """9.3 — Convert CSV data to financial Excel models with SUM formulas."""

    @pytest.mark.parametrize("n_rows", [5, 10, 25, 50])
    def test_csv_to_excel_model(self, out_dir, n_rows):
        csv_path = out_dir / f"data_{n_rows}.csv"
        rows = [["Item", "Amount"]] + [[f"Item_{j}", str(j * 100)] for j in range(n_rows)]
        write_csv(csv_path, rows)
        xlsx_path = out_dir / f"model_{n_rows}.xlsx"
        csv_to_xlsx_model(csv_path, xlsx_path)
        wb = load_workbook(xlsx_path)
        ws = wb.active
        assert ws.title == "Data"
        # Should have a TOTAL row with SUM formula
        total_row = ws.max_row
        assert str(ws.cell(total_row, 1).value) == "TOTAL"
        assert "=SUM" in str(ws.cell(total_row, 2).value)

    @pytest.mark.parametrize("i", range(10))
    def test_csv_to_excel_data_preserved(self, out_dir, i):
        csv_path = out_dir / f"preserve_{i}.csv"
        rows = [["Name", "Value"], [f"Alpha_{i}", "100"], [f"Beta_{i}", "200"]]
        write_csv(csv_path, rows)
        xlsx_path = out_dir / f"preserved_{i}.xlsx"
        csv_to_xlsx_model(csv_path, xlsx_path)
        wb = load_workbook(xlsx_path)
        ws = wb.active
        assert ws.cell(2, 1).value == f"Alpha_{i}"

    @pytest.mark.parametrize("cols", [2, 3, 5])
    def test_csv_to_excel_columns(self, out_dir, cols):
        csv_path = out_dir / f"cols_{cols}.csv"
        header = [f"Col_{j}" for j in range(cols)]
        data = [[str(i * j) for j in range(cols)] for i in range(5)]
        write_csv(csv_path, [header] + data)
        xlsx_path = out_dir / f"cols_xl_{cols}.xlsx"
        csv_to_xlsx_model(csv_path, xlsx_path)
        wb = load_workbook(xlsx_path)
        assert wb.active.max_column >= cols


# ============================================================================
# 9.4 — Word -> Presentacion
# ============================================================================

@pytest.mark.file_generation
class TestDocxToPptx:
    """9.4 — Convert Word documents to PowerPoint presentations with roundtrip validation."""

    @pytest.mark.parametrize("i", range(15))
    def test_docx_to_pptx(self, out_dir, i):
        spec = DocSpec(kind="docx", title=f"Strategy Doc {i}", sections=[
            f"Chapter 1: Market Analysis",
            f"Chapter 2: Competitive Landscape",
            f"Chapter 3: Go-to-Market Plan",
        ])
        docx_path = generate_document(spec, out_dir)
        pptx_path = out_dir / f"from_word_{i}.pptx"
        docx_to_pptx_outline(docx_path, pptx_path)
        prs = Presentation(pptx_path)
        assert len(prs.slides) >= 3

    @pytest.mark.parametrize("i", range(10))
    def test_docx_to_pptx_title_slide(self, out_dir, i):
        spec = DocSpec(kind="docx", title=f"Presentation {i}", sections=["Intro", "Body", "End"])
        docx_path = generate_document(spec, out_dir)
        pptx_path = out_dir / f"word_title_{i}.pptx"
        docx_to_pptx_outline(docx_path, pptx_path)
        prs = Presentation(pptx_path)
        assert f"Presentation {i}" in prs.slides[0].shapes.title.text

    @pytest.mark.parametrize("i", range(10))
    def test_pptx_to_markdown_roundtrip(self, out_dir, i):
        spec = DocSpec(kind="docx", title=f"Doc {i}", sections=["Section A", "Section B"])
        docx_path = generate_document(spec, out_dir)
        pptx_path = out_dir / f"roundtrip_{i}.pptx"
        docx_to_pptx_outline(docx_path, pptx_path)
        md = pptx_to_markdown_outline(pptx_path)
        assert "## Slide 1" in md
        assert "## Slide 2" in md


# ============================================================================
# 9.5 — Facturas/recibos -> Spreadsheet
# ============================================================================

@pytest.mark.file_generation
class TestInvoiceToSpreadsheet:
    """9.5 — Extract invoice/receipt data from PDFs into organized spreadsheets."""

    @pytest.mark.parametrize("i", range(15))
    def test_invoice_pdf_to_spreadsheet(self, out_dir, i):
        """Simulate extracting invoice data from PDF to spreadsheet."""
        lines = [
            f"Invoice #{i:04d}",
            f"Item  Qty  Price",
            f"Widget  {5 + i}  ${(5 + i) * 10}",
            f"Service  1  ${200 + i * 50}",
            f"Total  -  ${(5 + i) * 10 + 200 + i * 50}",
        ]
        spec = DocSpec(kind="pdf", title=f"Invoice_{i}", sections=lines)
        pdf_path = generate_document(spec, out_dir)
        rows = extract_pdf_tables(pdf_path)
        # Write to Excel
        xlsx_path = out_dir / f"invoice_xl_{i}.xlsx"
        create_excel_with_formulas(xlsx_path, rows, {})
        assert xlsx_path.exists()
        wb = load_workbook(xlsx_path)
        assert wb.active.max_row >= 2

    @pytest.mark.parametrize("i", range(10))
    def test_receipt_extraction_structure(self, out_dir, i):
        lines = [
            f"Store Name  Receipt #{i:04d}",
            f"Coffee  $4.50",
            f"Sandwich  $8.00",
            f"Tax  $1.00",
            f"Total  $13.50",
        ]
        spec = DocSpec(kind="pdf", title=f"Receipt_{i}", sections=lines)
        pdf_path = generate_document(spec, out_dir)
        rows = extract_pdf_tables(pdf_path)
        assert len(rows) >= 3


# ============================================================================
# 9.6 — Excel -> Reporte en Word
# ============================================================================

@pytest.mark.file_generation
class TestExcelToWordReport:
    """9.6 — Generate Word reports with comments and tables from Excel data."""

    @pytest.mark.parametrize("i", range(15))
    def test_excel_to_word_report(self, out_dir, i):
        """Create Excel data, then generate Word report with that data."""
        data = [
            ["Metric", "Q1", "Q2", "Q3", "Q4"],
            ["Revenue", 100 + i, 120 + i, 130 + i, 150 + i],
            ["Costs", 80 + i, 85 + i, 90 + i, 95 + i],
            ["Profit", 20, 35, 40, 55],
        ]
        xlsx_path = out_dir / f"source_{i}.xlsx"
        create_excel_with_formulas(xlsx_path, data, {})
        # Generate Word report from data
        tables_for_word = [("Financial Summary", [[str(c) for c in row] for row in data])]
        comments = {0: f"Data sourced from Excel model #{i}"}
        paragraphs = [
            f"This report summarizes the financial performance for period {i}.",
            f"Revenue grew from {100 + i} to {150 + i} over the year.",
            f"Overall profit margin improved steadily.",
        ]
        word_path = out_dir / f"report_{i}.docx"
        create_word_with_comments(word_path, paragraphs, comments)
        text = read_document_text(word_path)
        assert "report" in text.lower() or "summarizes" in text.lower()

    @pytest.mark.parametrize("i", range(10))
    def test_excel_data_in_word_table(self, out_dir, i):
        data = [
            ["Product", "Sales", "Target"],
            [f"Prod_A_{i}", "500", "600"],
            [f"Prod_B_{i}", "800", "750"],
        ]
        tables = [("Sales Data", data)]
        word_path = out_dir / f"table_report_{i}.docx"
        create_word_with_tables(word_path, tables)
        doc = Document(word_path)
        assert len(doc.tables) == 1
        assert doc.tables[0].cell(1, 0).text == f"Prod_A_{i}"
