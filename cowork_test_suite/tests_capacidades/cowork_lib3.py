"""
cowork_lib3
===========

Extended capability library for ILIAGPT.IO comprehensive test suite.
Adds deterministic implementations for capabilities not covered by
cowork_lib.py and cowork_lib2.py:

 1. Advanced Excel (pivot tables, conditional formatting rules, chart specs,
    dashboards, financial scenario models, budget trackers, data cleaning)
 2. Advanced PowerPoint (speaker notes, watermarks, GIF placeholders,
    multi-layout decks, notes-to-slides conversion)
 3. Advanced Word (redlines/comments, formatted tables, heading hierarchy,
    technical papers, legal documents)
 4. Advanced PDF (form fields, merge, split, table extraction)
 5. React/JSX/TSX generation
 6. Enhanced connectors (FactSet, DocuSign, Zoom transcription)
 7. Computer use extensions (spreadsheet filling, form completion)
 8. Scheduling extensions (cadence definitions, digest generation)
 9. Workspace extensions (per-folder instructions, live updates)
10. Availability & platform checks

Everything is deterministic and runs offline — no LLM or network access.
"""
from __future__ import annotations

import copy
import csv
import hashlib
import io
import json
import math
import os
import pathlib
import re
import shutil
import statistics
import textwrap
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional, Sequence, Tuple

# ---------------------------------------------------------------------------
# 1. Advanced Excel capabilities
# ---------------------------------------------------------------------------

from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.formatting.rule import CellIsRule, ColorScaleRule, FormulaRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter


def create_excel_with_formulas(
    out: pathlib.Path,
    rows: Sequence[Sequence[Any]],
    formulas: Dict[str, str],
) -> pathlib.Path:
    """Create an xlsx with data rows and named-cell formulas."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    for r in rows:
        ws.append(list(r))
    for cell_ref, formula in formulas.items():
        ws[cell_ref] = formula
    wb.save(out)
    return out


def create_excel_conditional_format(
    out: pathlib.Path,
    data: Sequence[Sequence[float]],
    threshold_high: float = 80,
    threshold_low: float = 20,
) -> pathlib.Path:
    """Create xlsx with conditional formatting rules."""
    wb = Workbook()
    ws = wb.active
    ws.title = "ConditionalData"
    for row in data:
        ws.append(list(row))
    max_row = ws.max_row
    max_col = ws.max_column
    cell_range = f"A1:{get_column_letter(max_col)}{max_row}"
    red_fill = PatternFill(start_color="FF9999", end_color="FF9999", fill_type="solid")
    green_fill = PatternFill(start_color="99FF99", end_color="99FF99", fill_type="solid")
    yellow_fill = PatternFill(start_color="FFFF99", end_color="FFFF99", fill_type="solid")
    ws.conditional_formatting.add(
        cell_range,
        CellIsRule(operator="greaterThan", formula=[str(threshold_high)], fill=green_fill),
    )
    ws.conditional_formatting.add(
        cell_range,
        CellIsRule(operator="lessThan", formula=[str(threshold_low)], fill=red_fill),
    )
    ws.conditional_formatting.add(
        cell_range,
        CellIsRule(
            operator="between",
            formula=[str(threshold_low), str(threshold_high)],
            fill=yellow_fill,
        ),
    )
    wb.save(out)
    return out


def create_excel_multi_sheet(
    out: pathlib.Path,
    sheets: Dict[str, Sequence[Sequence[Any]]],
) -> pathlib.Path:
    """Create xlsx with multiple named sheets."""
    wb = Workbook()
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name[:31])
        for r in rows:
            ws.append(list(r))
    wb.save(out)
    return out


def create_excel_with_chart(
    out: pathlib.Path,
    categories: Sequence[str],
    series_data: Dict[str, Sequence[float]],
    chart_type: str = "bar",
) -> pathlib.Path:
    """Create xlsx with embedded chart (bar, line, or pie)."""
    wb = Workbook()
    ws = wb.active
    ws.title = "ChartData"
    headers = ["Category"] + list(series_data.keys())
    ws.append(headers)
    for i, cat in enumerate(categories):
        row = [cat] + [series_data[s][i] for s in series_data]
        ws.append(row)
    chart_cls = {"bar": BarChart, "line": LineChart, "pie": PieChart}.get(chart_type, BarChart)
    chart = chart_cls()
    chart.title = "Generated Chart"
    chart.style = 10
    data_ref = Reference(ws, min_col=2, max_col=len(headers), min_row=1, max_row=len(categories) + 1)
    cats_ref = Reference(ws, min_col=1, min_row=2, max_row=len(categories) + 1)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)
    ws.add_chart(chart, "A" + str(len(categories) + 4))
    wb.save(out)
    return out


def create_financial_model(
    out: pathlib.Path,
    base_revenue: float,
    growth_rates: Sequence[float],
    cost_pct: float = 0.65,
    tax_rate: float = 0.21,
) -> pathlib.Path:
    """Create a multi-scenario financial model xlsx."""
    wb = Workbook()
    ws = wb.active
    ws.title = "FinancialModel"
    headers = ["Year", "Revenue", "COGS", "Gross Profit", "Tax", "Net Income"]
    ws.append(headers)
    for i, gr in enumerate(growth_rates):
        year = 2026 + i
        rev = base_revenue * (1 + gr) ** (i + 1)
        cogs = rev * cost_pct
        gross = rev - cogs
        tax = gross * tax_rate
        net = gross - tax
        ws.append([year, round(rev, 2), round(cogs, 2), round(gross, 2), round(tax, 2), round(net, 2)])
    # Scenario sheet
    ws2 = wb.create_sheet("Scenarios")
    ws2.append(["Scenario", "Growth Rate", "5Y Net Income"])
    scenarios = [("Bear", -0.05), ("Base", 0.05), ("Bull", 0.15)]
    for name, rate in scenarios:
        cumulative = sum(
            (base_revenue * (1 + rate) ** (y + 1)) * (1 - cost_pct) * (1 - tax_rate)
            for y in range(5)
        )
        ws2.append([name, rate, round(cumulative, 2)])
    wb.save(out)
    return out


def create_budget_tracker(
    out: pathlib.Path,
    categories: Sequence[str],
    budgets: Sequence[float],
    actuals: Sequence[float],
) -> pathlib.Path:
    """Create budget tracker with auto-calculated variances."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Category", "Budget", "Actual", "Variance", "% Variance", "Status"])
    for i, cat in enumerate(categories):
        row_num = i + 2
        ws.append([
            cat, budgets[i], actuals[i],
            f"=B{row_num}-C{row_num}",
            f"=IF(B{row_num}=0,0,D{row_num}/B{row_num})",
            f'=IF(D{row_num}>=0,"Under","Over")',
        ])
    total_row = len(categories) + 2
    ws.append([
        "TOTAL",
        f"=SUM(B2:B{total_row - 1})",
        f"=SUM(C2:C{total_row - 1})",
        f"=SUM(D2:D{total_row - 1})",
        f"=IF(B{total_row}=0,0,D{total_row}/B{total_row})",
        "",
    ])
    wb.save(out)
    return out


def create_pivot_table_sim(
    out: pathlib.Path,
    raw_data: Sequence[Dict[str, Any]],
    row_field: str,
    col_field: str,
    value_field: str,
) -> pathlib.Path:
    """Simulate a pivot table in xlsx (cross-tab aggregation)."""
    wb = Workbook()
    ws_raw = wb.active
    ws_raw.title = "RawData"
    if not raw_data:
        wb.save(out)
        return out
    headers = list(raw_data[0].keys())
    ws_raw.append(headers)
    for rec in raw_data:
        ws_raw.append([rec.get(h) for h in headers])
    # Build pivot
    row_vals = sorted(set(r[row_field] for r in raw_data))
    col_vals = sorted(set(r[col_field] for r in raw_data))
    pivot: Dict[str, Dict[str, float]] = {rv: {cv: 0.0 for cv in col_vals} for rv in row_vals}
    for rec in raw_data:
        pivot[rec[row_field]][rec[col_field]] += float(rec[value_field])
    ws_pivot = wb.create_sheet("Pivot")
    ws_pivot.append([row_field] + col_vals + ["Total"])
    for rv in row_vals:
        row = [rv] + [pivot[rv][cv] for cv in col_vals]
        row.append(sum(pivot[rv].values()))
        ws_pivot.append(row)
    wb.save(out)
    return out


def clean_dataset(
    data: Sequence[Dict[str, Any]],
    *,
    drop_nulls: bool = True,
    dedup_key: Optional[str] = None,
    strip_whitespace: bool = True,
) -> List[Dict[str, Any]]:
    """Deterministic data cleaning pipeline."""
    result: List[Dict[str, Any]] = []
    seen_keys: set = set()
    for rec in data:
        row = dict(rec)
        if strip_whitespace:
            row = {k: v.strip() if isinstance(v, str) else v for k, v in row.items()}
        if drop_nulls and any(v is None or v == "" for v in row.values()):
            continue
        if dedup_key and row.get(dedup_key) in seen_keys:
            continue
        if dedup_key:
            seen_keys.add(row[dedup_key])
        result.append(row)
    return result


def create_dashboard_xlsx(
    out: pathlib.Path,
    kpis: Dict[str, float],
    monthly_data: Dict[str, Sequence[float]],
) -> pathlib.Path:
    """Create a dashboard xlsx with KPIs and trend charts."""
    wb = Workbook()
    ws_kpi = wb.active
    ws_kpi.title = "KPI_Dashboard"
    ws_kpi.append(["KPI", "Value", "Target", "Status"])
    for name, val in kpis.items():
        target = val * 1.1
        status = "On Track" if val >= target * 0.9 else "At Risk"
        ws_kpi.append([name, val, round(target, 2), status])
    ws_trend = wb.create_sheet("Trends")
    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    ws_trend.append(["Month"] + list(monthly_data.keys()))
    for i, m in enumerate(months):
        row = [m] + [monthly_data[k][i] if i < len(monthly_data[k]) else 0 for k in monthly_data]
        ws_trend.append(row)
    chart = LineChart()
    chart.title = "Monthly Trends"
    data_ref = Reference(ws_trend, min_col=2, max_col=1 + len(monthly_data), min_row=1, max_row=13)
    cats_ref = Reference(ws_trend, min_col=1, min_row=2, max_row=13)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)
    ws_trend.add_chart(chart, "A15")
    wb.save(out)
    return out


# ---------------------------------------------------------------------------
# 2. Advanced PowerPoint capabilities
# ---------------------------------------------------------------------------

from pptx import Presentation
from pptx.util import Inches, Pt, Emu


def create_pptx_with_notes(
    out: pathlib.Path,
    slides: Sequence[Dict[str, str]],
) -> pathlib.Path:
    """Create pptx where each slide has title, body, and speaker notes."""
    prs = Presentation()
    for s in slides:
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s.get("title", "")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = s.get("body", "")
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = s.get("notes", "")
    prs.save(out)
    return out


def create_pptx_with_watermark(
    out: pathlib.Path,
    slides_text: Sequence[str],
    watermark: str = "CONFIDENTIAL",
) -> pathlib.Path:
    """Create pptx with a text watermark on every slide."""
    prs = Presentation()
    for txt in slides_text:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = txt
        from pptx.util import Pt as PtU
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = watermark
        p.font.size = Pt(48)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    prs.save(out)
    return out


def create_pptx_from_doc_outline(
    doc_paragraphs: Sequence[str],
    out: pathlib.Path,
) -> pathlib.Path:
    """Convert document paragraphs into a presentation outline."""
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = doc_paragraphs[0] if doc_paragraphs else "Untitled"
    for para in doc_paragraphs[1:]:
        if not para.strip():
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = para[:60]
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = para
    prs.save(out)
    return out


def create_pptx_multi_layout(
    out: pathlib.Path,
    content: Sequence[Dict[str, Any]],
) -> pathlib.Path:
    """Create pptx using different slide layouts."""
    prs = Presentation()
    for item in content:
        layout_idx = min(item.get("layout", 0), len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        if slide.shapes.title:
            slide.shapes.title.text = item.get("title", "")
        notes = slide.notes_slide
        notes.notes_text_frame.text = item.get("notes", "")
    prs.save(out)
    return out


# ---------------------------------------------------------------------------
# 3. Advanced Word capabilities
# ---------------------------------------------------------------------------

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGBColor


def create_word_with_hierarchy(
    out: pathlib.Path,
    structure: Sequence[Tuple[int, str, str]],
) -> pathlib.Path:
    """Create docx with heading hierarchy: [(level, heading, body), ...]."""
    doc = Document()
    for level, heading, body in structure:
        doc.add_heading(heading, level=min(level, 9))
        if body:
            doc.add_paragraph(body)
    doc.save(out)
    return out


def create_word_with_tables(
    out: pathlib.Path,
    tables_data: Sequence[Tuple[str, Sequence[Sequence[str]]]],
) -> pathlib.Path:
    """Create docx with formatted tables: [(table_name, rows), ...]."""
    doc = Document()
    for name, rows in tables_data:
        doc.add_heading(name, level=2)
        if not rows:
            continue
        table = doc.add_table(rows=len(rows), cols=len(rows[0]))
        table.style = "Table Grid"
        for i, row in enumerate(rows):
            for j, cell_val in enumerate(row):
                table.cell(i, j).text = str(cell_val)
                if i == 0:
                    for paragraph in table.cell(i, j).paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
        doc.add_paragraph("")
    doc.save(out)
    return out


def create_word_with_comments(
    out: pathlib.Path,
    paragraphs: Sequence[str],
    comments: Dict[int, str],
) -> pathlib.Path:
    """Create docx with comment annotations (simulated via inline markup).

    Real OOXML comments require low-level XML manipulation.
    We simulate by adding highlighted [COMMENT: ...] markers.
    """
    doc = Document()
    doc.add_heading("Document with Review Comments", level=1)
    for i, para_text in enumerate(paragraphs):
        p = doc.add_paragraph()
        run = p.add_run(para_text)
        if i in comments:
            comment_run = p.add_run(f"  [COMMENT: {comments[i]}]")
            comment_run.font.color.rgb = DocRGBColor(0xFF, 0x00, 0x00)
            comment_run.font.italic = True
            comment_run.font.size = DocPt(9)
    doc.save(out)
    return out


def create_word_redline(
    out: pathlib.Path,
    original: Sequence[str],
    revised: Sequence[str],
) -> pathlib.Path:
    """Simulate a redline document showing original (strikethrough) vs revised."""
    doc = Document()
    doc.add_heading("Redline Comparison", level=1)
    for orig, rev in zip(original, revised):
        p = doc.add_paragraph()
        if orig != rev:
            del_run = p.add_run(orig)
            del_run.font.strike = True
            del_run.font.color.rgb = DocRGBColor(0xFF, 0x00, 0x00)
            p.add_run("  →  ")
            add_run = p.add_run(rev)
            add_run.font.color.rgb = DocRGBColor(0x00, 0x80, 0x00)
            add_run.font.underline = True
        else:
            p.add_run(orig)
    doc.save(out)
    return out


def create_technical_paper(
    out: pathlib.Path,
    title: str,
    abstract: str,
    sections: Sequence[Tuple[str, str]],
    references: Sequence[str],
) -> pathlib.Path:
    """Create a technical paper in docx format."""
    doc = Document()
    doc.add_heading(title, level=0)
    doc.add_heading("Abstract", level=1)
    doc.add_paragraph(abstract).italic = True
    for heading, content in sections:
        doc.add_heading(heading, level=1)
        doc.add_paragraph(content)
    if references:
        doc.add_heading("References", level=1)
        for i, ref in enumerate(references, 1):
            doc.add_paragraph(f"[{i}] {ref}")
    doc.save(out)
    return out


# ---------------------------------------------------------------------------
# 4. Advanced PDF capabilities
# ---------------------------------------------------------------------------

from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.lib.pagesizes import letter
from pypdf import PdfReader, PdfWriter


def create_pdf_with_fields(
    out: pathlib.Path,
    field_specs: Dict[str, str],
) -> pathlib.Path:
    """Create a PDF and simulate form fields via text annotations."""
    c = pdf_canvas.Canvas(str(out), pagesize=letter)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 720, "Form Document")
    c.setFont("Helvetica", 11)
    y = 680
    for field_name, default_val in field_specs.items():
        c.drawString(72, y, f"{field_name}: ")
        c.drawString(200, y, f"[{default_val}]")
        y -= 24
    c.save()
    return out


def fill_pdf_form_sim(
    template_path: pathlib.Path,
    values: Dict[str, str],
    out: pathlib.Path,
) -> pathlib.Path:
    """Simulate filling a PDF form by creating a new PDF with filled values."""
    c = pdf_canvas.Canvas(str(out), pagesize=letter)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 720, "Filled Form")
    c.setFont("Helvetica", 11)
    y = 680
    for field_name, value in values.items():
        c.drawString(72, y, f"{field_name}: {value}")
        y -= 24
    c.save()
    return out


def merge_pdfs(paths: Sequence[pathlib.Path], out: pathlib.Path) -> pathlib.Path:
    """Merge multiple PDFs into one."""
    writer = PdfWriter()
    for p in paths:
        reader = PdfReader(str(p))
        for page in reader.pages:
            writer.add_page(page)
    with open(out, "wb") as f:
        writer.write(f)
    return out


def split_pdf(
    src: pathlib.Path,
    out_dir: pathlib.Path,
) -> List[pathlib.Path]:
    """Split a PDF into individual page files."""
    reader = PdfReader(str(src))
    out_dir.mkdir(parents=True, exist_ok=True)
    results: List[pathlib.Path] = []
    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)
        page_path = out_dir / f"page_{i + 1}.pdf"
        with open(page_path, "wb") as f:
            writer.write(f)
        results.append(page_path)
    return results


def extract_pdf_tables(
    src: pathlib.Path,
) -> List[List[str]]:
    """Extract text lines from PDF (simulated table extraction)."""
    reader = PdfReader(str(src))
    rows: List[List[str]] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        for line in text.strip().split("\n"):
            cells = re.split(r"\s{2,}", line.strip())
            if cells and any(c.strip() for c in cells):
                rows.append(cells)
    return rows


# ---------------------------------------------------------------------------
# 5. React / JSX / TSX generation
# ---------------------------------------------------------------------------

def generate_react_component(
    name: str,
    props: Dict[str, str],
    body_jsx: str = "<div>Hello</div>",
    typescript: bool = True,
) -> str:
    """Generate a React functional component."""
    ext = "tsx" if typescript else "jsx"
    if typescript and props:
        interface_fields = "\n".join(f"  {k}: {v};" for k, v in props.items())
        interface = f"interface {name}Props {{\n{interface_fields}\n}}\n\n"
        args = f"{{ {', '.join(props.keys())} }}: {name}Props"
    else:
        interface = ""
        args = f"{{ {', '.join(props.keys())} }}" if props else ""
    component = textwrap.dedent(f"""\
        {interface}export default function {name}({args}) {{
          return (
            {body_jsx}
          );
        }}
    """)
    return component


# ---------------------------------------------------------------------------
# 6. Enhanced connectors
# ---------------------------------------------------------------------------

@dataclass
class ZoomTranscript:
    meeting_id: str
    duration_min: int
    participants: List[str]
    segments: List[Dict[str, str]]  # {"speaker": ..., "text": ...}

    def summary(self) -> str:
        speakers = set(s["speaker"] for s in self.segments)
        return (
            f"Meeting {self.meeting_id}: {self.duration_min}min, "
            f"{len(speakers)} speakers, {len(self.segments)} segments"
        )

    def action_items(self) -> List[str]:
        items: List[str] = []
        for seg in self.segments:
            text = seg["text"].lower()
            if any(kw in text for kw in ["action item", "todo", "will do", "follow up", "next step"]):
                items.append(f"{seg['speaker']}: {seg['text']}")
        return items


@dataclass
class DocuSignEnvelope:
    envelope_id: str
    signers: List[str]
    status: str = "created"
    signed_by: List[str] = field(default_factory=list)

    def send(self) -> str:
        self.status = "sent"
        return f"sent:{self.envelope_id}"

    def sign(self, signer: str) -> bool:
        if signer not in self.signers or signer in self.signed_by:
            return False
        self.signed_by.append(signer)
        if set(self.signed_by) == set(self.signers):
            self.status = "completed"
        return True


@dataclass
class FactSetQuery:
    ticker: str
    metrics: List[str]

    def execute(self) -> Dict[str, Any]:
        """Deterministic mock financial data."""
        h = int(hashlib.md5(self.ticker.encode()).hexdigest()[:8], 16)
        result: Dict[str, Any] = {"ticker": self.ticker}
        for metric in self.metrics:
            if metric == "revenue":
                result[metric] = round((h % 10000) * 1e6, 2)
            elif metric == "pe_ratio":
                result[metric] = round((h % 50) + 5 + (h % 100) / 100, 2)
            elif metric == "market_cap":
                result[metric] = round((h % 500) * 1e9, 2)
            else:
                result[metric] = round((h % 1000) / 10, 2)
        return result


# ---------------------------------------------------------------------------
# 7. Computer use extensions
# ---------------------------------------------------------------------------

@dataclass
class SpreadsheetApp:
    """Simulates filling a spreadsheet application."""
    cells: Dict[Tuple[int, int], Any] = field(default_factory=dict)

    def set_cell(self, row: int, col: int, value: Any) -> None:
        self.cells[(row, col)] = value

    def get_cell(self, row: int, col: int) -> Any:
        return self.cells.get((row, col))

    def fill_range(self, start_row: int, start_col: int, data: Sequence[Sequence[Any]]) -> int:
        count = 0
        for i, row in enumerate(data):
            for j, val in enumerate(row):
                self.set_cell(start_row + i, start_col + j, val)
                count += 1
        return count


@dataclass
class WebFormApp:
    """Simulates completing a web form."""
    fields: Dict[str, str] = field(default_factory=dict)
    submitted: bool = False
    validation_rules: Dict[str, str] = field(default_factory=dict)

    def fill(self, field_name: str, value: str) -> None:
        self.fields[field_name] = value

    def validate(self) -> List[str]:
        errors: List[str] = []
        for fname, rule in self.validation_rules.items():
            val = self.fields.get(fname, "")
            if rule == "required" and not val:
                errors.append(f"{fname} is required")
            elif rule == "email" and "@" not in val:
                errors.append(f"{fname} must be a valid email")
            elif rule == "numeric" and not val.replace(".", "").isdigit():
                errors.append(f"{fname} must be numeric")
        return errors

    def submit(self) -> Tuple[bool, List[str]]:
        errors = self.validate()
        if not errors:
            self.submitted = True
        return (not errors, errors)


# ---------------------------------------------------------------------------
# 8. Scheduling extensions
# ---------------------------------------------------------------------------

@dataclass
class ScheduledTask:
    name: str
    cron_expr: str
    action: str
    enabled: bool = True
    last_run: Optional[datetime] = None
    run_count: int = 0

    def execute(self, now: datetime) -> Dict[str, Any]:
        self.last_run = now
        self.run_count += 1
        return {"task": self.name, "action": self.action, "at": now.isoformat(), "run": self.run_count}


class TaskScheduler:
    def __init__(self) -> None:
        self.tasks: List[ScheduledTask] = []

    def add(self, name: str, cron: str, action: str) -> ScheduledTask:
        t = ScheduledTask(name=name, cron_expr=cron, action=action)
        self.tasks.append(t)
        return t

    def list_tasks(self) -> List[Dict[str, Any]]:
        return [
            {"name": t.name, "cron": t.cron_expr, "enabled": t.enabled, "runs": t.run_count}
            for t in self.tasks
        ]

    def toggle(self, name: str) -> bool:
        for t in self.tasks:
            if t.name == name:
                t.enabled = not t.enabled
                return t.enabled
        raise KeyError(name)

    def run_due(self, now: datetime) -> List[Dict[str, Any]]:
        results: List[Dict[str, Any]] = []
        for t in self.tasks:
            if t.enabled:
                results.append(t.execute(now))
        return results


def generate_digest(
    source: str,
    items: Sequence[Dict[str, str]],
) -> str:
    """Generate a digest summary from items."""
    lines = [f"# {source} Digest", f"Generated: {datetime.now().date()}", ""]
    for item in items:
        lines.append(f"- **{item.get('title', 'Untitled')}**: {item.get('summary', 'N/A')}")
    lines.append(f"\nTotal items: {len(items)}")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# 9. Workspace extensions
# ---------------------------------------------------------------------------

@dataclass
class FolderInstructions:
    path: pathlib.Path
    instructions: Dict[str, str] = field(default_factory=dict)

    def set(self, key: str, value: str) -> None:
        self.instructions[key] = value
        self._save()

    def get(self, key: str) -> Optional[str]:
        return self.instructions.get(key)

    def update_live(self, key: str, value: str) -> None:
        self.instructions[key] = value
        self._save()

    def _save(self) -> None:
        self.path.mkdir(parents=True, exist_ok=True)
        (self.path / ".instructions.json").write_text(
            json.dumps(self.instructions, indent=2)
        )

    @classmethod
    def load(cls, path: pathlib.Path) -> "FolderInstructions":
        fp = path / ".instructions.json"
        if fp.exists():
            data = json.loads(fp.read_text())
        else:
            data = {}
        return cls(path=path, instructions=data)


@dataclass
class GlobalInstructions:
    tone: str = "professional"
    format: str = "concise"
    role_context: str = ""
    custom_rules: List[str] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "tone": self.tone,
            "format": self.format,
            "role_context": self.role_context,
            "custom_rules": self.custom_rules,
        }

    def apply_to_prompt(self, prompt: str) -> str:
        prefix = f"[Tone: {self.tone}] [Format: {self.format}]"
        if self.role_context:
            prefix += f" [Role: {self.role_context}]"
        return f"{prefix}\n{prompt}"


# ---------------------------------------------------------------------------
# 10. Availability & platform checks
# ---------------------------------------------------------------------------

SUPPORTED_PLATFORMS = {"macos", "windows"}
SUPPORTED_PLANS = {"pro", "max", "team", "enterprise"}
MAX_FILE_SIZE_MB = 30
DISPATCH_PLANS = {"pro", "max"}
SAVE_DESTINATIONS = {"local", "google_drive"}


def check_platform(platform: str) -> bool:
    return platform.lower() in SUPPORTED_PLATFORMS


def check_plan_access(plan: str, feature: str) -> bool:
    plan = plan.lower()
    if plan not in SUPPORTED_PLANS:
        return False
    if feature == "dispatch" and plan not in DISPATCH_PLANS:
        return False
    if feature == "rbac" and plan != "enterprise":
        return False
    if feature == "private_marketplace" and plan != "enterprise":
        return False
    return True


def check_file_size(size_mb: float) -> bool:
    return size_mb <= MAX_FILE_SIZE_MB


# ---------------------------------------------------------------------------
# 11. Variance analysis helper
# ---------------------------------------------------------------------------

def variance_analysis(
    budget: Dict[str, float],
    actual: Dict[str, float],
) -> Dict[str, Dict[str, float]]:
    """Compute variance analysis for budget vs actual."""
    result: Dict[str, Dict[str, float]] = {}
    for key in budget:
        b = budget[key]
        a = actual.get(key, 0.0)
        var = a - b
        pct = (var / b * 100) if b != 0 else 0.0
        result[key] = {
            "budget": b,
            "actual": a,
            "variance": var,
            "variance_pct": round(pct, 2),
            "favorable": var <= 0,
        }
    return result


# ---------------------------------------------------------------------------
# 12. Time series analysis
# ---------------------------------------------------------------------------

def detect_trend(series: Sequence[float]) -> str:
    """Detect if series is trending up, down, or flat."""
    if len(series) < 2:
        return "insufficient_data"
    n = len(series)
    xs = list(range(n))
    mx = sum(xs) / n
    my = sum(series) / n
    num = sum((x - mx) * (y - my) for x, y in zip(xs, series))
    den = sum((x - mx) ** 2 for x in xs)
    slope = num / den if den else 0
    if abs(slope) < 0.01 * (max(series) - min(series) + 1e-9):
        return "flat"
    return "up" if slope > 0 else "down"


def seasonal_decompose(
    series: Sequence[float],
    period: int = 4,
) -> Dict[str, List[float]]:
    """Simple additive seasonal decomposition."""
    n = len(series)
    trend: List[float] = []
    for i in range(n):
        start = max(0, i - period // 2)
        end = min(n, i + period // 2 + 1)
        trend.append(sum(series[start:end]) / (end - start))
    seasonal: List[float] = []
    for i in range(n):
        seasonal.append(series[i] - trend[i])
    residual: List[float] = [0.0] * n
    return {"trend": trend, "seasonal": seasonal, "residual": residual}
