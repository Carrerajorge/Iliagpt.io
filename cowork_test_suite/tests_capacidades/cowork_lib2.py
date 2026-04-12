"""
cowork_lib2
===========

Extension library for IliaGPT / Cowork capability tests. Adds deterministic,
pure-Python implementations of every capability category the product claims:

 1. File generation extras (md, html, csv/tsv, json, latex, png chart, code files)
 2. File management (rename, dedupe, organize by extension/content hash)
 3. Data science (stats, outliers, forecasting, tiny ML)
 4. Synthesis / research (multi-doc reading + citations + contradiction hints)
 5. Format conversion (csv->xlsx model, md->html, md->docx text, word->pptx text,
    pptx->md outline, pdf text extraction -> xlsx table)
 6. Browser automation mock (URL parsing, form filling, DOM-ish operations)
 7. Computer use mock (command parsing + safe echo)
 8. Scheduled tasks (cron parser + next-fire)
 9. Dispatch mock (mobile -> desktop message queue)
10. Connector mock (Google Drive / Gmail / Slack / Jira / Notion / Github stubs)
11. Plugins / skill-creator mock
12. Code execution (python + node subprocess with sandbox timeout)
13. Sub-agents / task decomposition
14. Project workspaces (persistent state)
15. Security / governance (path sandbox, deletion guard, egress allowlist)
16. Enterprise (RBAC, budget, telemetry aggregation)
17. Domain templates (legal / finance / marketing / ops / hr / research)

Everything is deterministic and runs offline.
"""
from __future__ import annotations

import csv
import dataclasses
import hashlib
import io
import json
import math
import os
import pathlib
import random
import re
import shutil
import statistics
import subprocess
import sys
import tempfile
import textwrap
import time
import urllib.parse
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timedelta
from typing import Any, Callable, Dict, Iterable, List, Optional, Sequence, Tuple

# ---------------------------------------------------------------------------
# 1. File generation extras
# ---------------------------------------------------------------------------

def write_markdown(path: pathlib.Path, title: str, sections: Sequence[str]) -> pathlib.Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    body = [f"# {title}", ""]
    for i, s in enumerate(sections, 1):
        body.append(f"## Section {i}")
        body.append(s)
        body.append("")
    path.write_text("\n".join(body), encoding="utf-8")
    return path


def write_html(path: pathlib.Path, title: str, body_html: str) -> pathlib.Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    html = (
        f"<!DOCTYPE html><html lang=\"es\"><head><meta charset=\"utf-8\">"
        f"<title>{title}</title></head><body>{body_html}</body></html>"
    )
    path.write_text(html, encoding="utf-8")
    return path


def write_csv(path: pathlib.Path, rows: Sequence[Sequence[Any]], delimiter: str = ",") -> pathlib.Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f, delimiter=delimiter)
        for row in rows:
            w.writerow(row)
    return path


def write_json(path: pathlib.Path, obj: Any) -> pathlib.Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def write_latex(path: pathlib.Path, title: str, author: str, body: str) -> pathlib.Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    tex = textwrap.dedent(rf"""
    \documentclass{{article}}
    \title{{{title}}}
    \author{{{author}}}
    \begin{{document}}
    \maketitle
    {body}
    \end{{document}}
    """).strip()
    path.write_text(tex, encoding="utf-8")
    return path


def write_png_chart(path: pathlib.Path, xs: Sequence[float], ys: Sequence[float], title: str = "chart") -> pathlib.Path:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    path.parent.mkdir(parents=True, exist_ok=True)
    fig, ax = plt.subplots(figsize=(4, 3), dpi=100)
    ax.plot(list(xs), list(ys), color="#2563eb", linewidth=2)
    ax.set_title(title)
    ax.grid(True, alpha=0.3)
    fig.tight_layout()
    fig.savefig(path)
    plt.close(fig)
    return path


def write_code_file(path: pathlib.Path, language: str, content: str) -> pathlib.Path:
    ext_map = {
        "python": "py", "node": "js", "javascript": "js", "typescript": "ts",
        "go": "go", "rust": "rs", "java": "java", "c": "c", "cpp": "cpp",
        "sh": "sh", "ruby": "rb", "php": "php",
    }
    ext = ext_map.get(language, language)
    path = path.with_suffix(f".{ext}") if not path.suffix else path
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


# ---------------------------------------------------------------------------
# 2. File management
# ---------------------------------------------------------------------------

def file_sha256(p: pathlib.Path) -> str:
    h = hashlib.sha256()
    with p.open("rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def dedupe_files(root: pathlib.Path) -> Tuple[int, int]:
    """Move duplicates (by content hash) into a 'duplicates/' subfolder.

    Returns (unique_count, duplicate_count).
    """
    seen: dict[str, pathlib.Path] = {}
    dup_dir = root / "duplicates"
    uniques = 0
    dups = 0
    for p in sorted(root.rglob("*")):
        if not p.is_file() or dup_dir in p.parents:
            continue
        digest = file_sha256(p)
        if digest in seen:
            dup_dir.mkdir(exist_ok=True)
            dst = dup_dir / p.name
            i = 1
            while dst.exists():
                dst = dup_dir / f"{p.stem}_{i}{p.suffix}"
                i += 1
            shutil.move(str(p), str(dst))
            dups += 1
        else:
            seen[digest] = p
            uniques += 1
    return uniques, dups


def rename_with_date_prefix(root: pathlib.Path, when: datetime) -> int:
    """Rename every top-level file in root prefixing with YYYY-MM-DD_."""
    prefix = when.strftime("%Y-%m-%d_")
    n = 0
    for p in sorted(root.iterdir()):
        if not p.is_file():
            continue
        if p.name.startswith(prefix):
            continue
        new_name = root / f"{prefix}{p.name}"
        p.rename(new_name)
        n += 1
    return n


def organize_by_extension(root: pathlib.Path) -> Dict[str, int]:
    """Move files into subfolders named by extension. Returns counts per ext."""
    counts: Dict[str, int] = {}
    for p in sorted(root.iterdir()):
        if not p.is_file():
            continue
        ext = (p.suffix[1:] if p.suffix else "noext").lower()
        sub = root / ext
        sub.mkdir(exist_ok=True)
        shutil.move(str(p), str(sub / p.name))
        counts[ext] = counts.get(ext, 0) + 1
    return counts


def safe_delete(path: pathlib.Path, *, allow: bool) -> bool:
    """Protection-against-deletion helper: only deletes if explicit allow=True."""
    if not allow:
        return False
    if not path.exists():
        return False
    if path.is_dir():
        shutil.rmtree(path)
    else:
        path.unlink()
    return True


# ---------------------------------------------------------------------------
# 3. Data science
# ---------------------------------------------------------------------------

def zscore_outliers(values: Sequence[float], threshold: float = 2.5) -> List[int]:
    if len(values) < 2:
        return []
    mu = statistics.fmean(values)
    sd = statistics.pstdev(values) or 1e-9
    return [i for i, v in enumerate(values) if abs((v - mu) / sd) > threshold]


def moving_average_forecast(series: Sequence[float], window: int = 3, horizon: int = 3) -> List[float]:
    s = list(series)
    out: List[float] = []
    for _ in range(horizon):
        w = s[-window:] if len(s) >= window else s
        nxt = sum(w) / len(w)
        out.append(nxt)
        s.append(nxt)
    return out


def linear_regression(xs: Sequence[float], ys: Sequence[float]) -> Tuple[float, float, float]:
    """Return (slope, intercept, r^2) using closed-form least squares."""
    n = len(xs)
    if n < 2:
        return 0.0, (ys[0] if ys else 0.0), 0.0
    mx = sum(xs) / n
    my = sum(ys) / n
    num = sum((x - mx) * (y - my) for x, y in zip(xs, ys))
    den = sum((x - mx) ** 2 for x in xs)
    if den == 0:
        return 0.0, my, 0.0
    slope = num / den
    intercept = my - slope * mx
    ss_tot = sum((y - my) ** 2 for y in ys)
    ss_res = sum((y - (slope * x + intercept)) ** 2 for x, y in zip(xs, ys))
    r2 = 1 - ss_res / ss_tot if ss_tot else 1.0
    return slope, intercept, r2


def knn_classify(point: Sequence[float], labeled: Sequence[Tuple[Sequence[float], Any]], k: int = 3) -> Any:
    dists = []
    for feats, lbl in labeled:
        d = math.sqrt(sum((a - b) ** 2 for a, b in zip(point, feats)))
        dists.append((d, lbl))
    dists.sort(key=lambda x: x[0])
    topk = [lbl for _, lbl in dists[:k]]
    # majority vote
    counts: Dict[Any, int] = {}
    for l in topk:
        counts[l] = counts.get(l, 0) + 1
    return max(counts.items(), key=lambda x: x[1])[0]


def cross_tab(rows: Sequence[Tuple[str, str]]) -> Dict[str, Dict[str, int]]:
    """rows is [(row_key, col_key), ...]. Returns nested dict of counts."""
    ct: Dict[str, Dict[str, int]] = {}
    for r, c in rows:
        ct.setdefault(r, {}).setdefault(c, 0)
        ct[r][c] += 1
    return ct


# ---------------------------------------------------------------------------
# 4. Synthesis / research
# ---------------------------------------------------------------------------

@dataclass
class Source:
    id: str
    title: str
    text: str


def extract_sentences(text: str) -> List[str]:
    parts = re.split(r"(?<=[\.\!\?])\s+", text.strip())
    return [p for p in parts if p]


def synthesize(sources: Sequence[Source], query: str, top_k: int = 3) -> Dict[str, Any]:
    """Return a synthesis with cited sentences from each source relevant to query."""
    q = set(re.findall(r"\w+", query.lower()))
    hits: List[Tuple[str, str, int]] = []  # (source_id, sentence, score)
    for src in sources:
        for sent in extract_sentences(src.text):
            words = set(re.findall(r"\w+", sent.lower()))
            score = len(q & words)
            if score > 0:
                hits.append((src.id, sent, score))
    hits.sort(key=lambda t: t[2], reverse=True)
    top = hits[:top_k]
    return {
        "query": query,
        "citations": [{"source": h[0], "sentence": h[1]} for h in top],
        "summary": " ".join(h[1] for h in top),
    }


def detect_contradictions(sources: Sequence[Source]) -> List[Tuple[str, str, str]]:
    """Very naive: detect 'X is Y' vs 'X is not Y' pairs across sources."""
    out: List[Tuple[str, str, str]] = []
    claims: Dict[str, Tuple[str, str]] = {}  # subject -> (source_id, sentence)
    for src in sources:
        for sent in extract_sentences(src.text):
            m = re.match(r"(?i)^(.*?)\s+(is|are|was|were)\s+(not\s+)?(.*?)[\.\!\?]?$", sent.strip())
            if not m:
                continue
            subj = m.group(1).strip().lower()
            negated = bool(m.group(3))
            pred = m.group(4).strip().lower()
            key = f"{subj}|{pred}"
            if key in claims:
                prev_src, prev_sent = claims[key]
                prev_negated = "not " in prev_sent.lower()
                if prev_negated != negated and prev_src != src.id:
                    out.append((src.id, prev_src, f"{subj} / {pred}"))
            claims[key] = (src.id, sent)
    return out


# ---------------------------------------------------------------------------
# 5. Format conversion
# ---------------------------------------------------------------------------

def markdown_to_html(md: str) -> str:
    import markdown2
    return markdown2.markdown(md, extras=["tables", "fenced-code-blocks"])


def csv_to_xlsx_model(csv_path: pathlib.Path, out_path: pathlib.Path) -> pathlib.Path:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    with csv_path.open(encoding="utf-8") as f:
        for row in csv.reader(f):
            ws.append(row)
    # Add a totals row with a SUM formula over numeric column B if present.
    if ws.max_row >= 2 and ws.max_column >= 2:
        total_row = ws.max_row + 1
        ws.cell(row=total_row, column=1, value="TOTAL")
        ws.cell(
            row=total_row, column=2,
            value=f"=SUM(B2:B{ws.max_row})",
        )
    wb.save(out_path)
    return out_path


def docx_to_pptx_outline(docx_path: pathlib.Path, out_path: pathlib.Path) -> pathlib.Path:
    from docx import Document
    from pptx import Presentation

    doc = Document(docx_path)
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = (doc.paragraphs[0].text if doc.paragraphs else "Outline")
    for para in doc.paragraphs[1:]:
        if not para.text.strip():
            continue
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = para.text[:60]
    prs.save(out_path)
    return out_path


def pptx_to_markdown_outline(pptx_path: pathlib.Path) -> str:
    from pptx import Presentation
    prs = Presentation(pptx_path)
    out: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for j, para in enumerate(shape.text_frame.paragraphs):
                    if j == 0 and not title:
                        title = para.text
                    else:
                        if para.text:
                            body.append(para.text)
        out.append(f"## Slide {i}: {title}")
        for b in body:
            out.append(f"- {b}")
    return "\n".join(out)


# ---------------------------------------------------------------------------
# 6. Browser automation mock
# ---------------------------------------------------------------------------

@dataclass
class MockBrowser:
    history: List[str] = field(default_factory=list)
    forms: Dict[str, Dict[str, str]] = field(default_factory=dict)

    def navigate(self, url: str) -> str:
        parsed = urllib.parse.urlparse(url)
        if parsed.scheme not in ("http", "https"):
            raise ValueError("unsafe scheme")
        self.history.append(url)
        return f"loaded:{parsed.netloc}{parsed.path}"

    def fill_form(self, form_id: str, fields: Dict[str, str]) -> None:
        self.forms[form_id] = dict(fields)

    def submit(self, form_id: str) -> Dict[str, str]:
        return self.forms.get(form_id, {})

    def screenshot(self, url: str) -> bytes:
        return hashlib.sha256(url.encode()).digest()

    def eval_js(self, expr: str) -> Any:
        # tiny deterministic sandbox: supports 1+2, Math.floor(x), "abc".length
        if expr.strip() == "1+2":
            return 3
        m = re.match(r"Math\.floor\(([-0-9\.]+)\)", expr.strip())
        if m:
            return math.floor(float(m.group(1)))
        m = re.match(r'"([^"]*)"\.length', expr.strip())
        if m:
            return len(m.group(1))
        raise NotImplementedError(expr)


# ---------------------------------------------------------------------------
# 7. Computer use mock
# ---------------------------------------------------------------------------

ALLOWED_APPS = {"Excel", "Word", "PowerPoint", "Chrome", "Safari", "Notes"}


def open_app(name: str) -> str:
    if name not in ALLOWED_APPS:
        raise PermissionError(f"App {name!r} not allowed")
    return f"opened:{name}"


# ---------------------------------------------------------------------------
# 8. Scheduled tasks (tiny cron)
# ---------------------------------------------------------------------------

def _parse_cron_field(field: str, lo: int, hi: int) -> List[int]:
    if field == "*":
        return list(range(lo, hi + 1))
    out: set[int] = set()
    for piece in field.split(","):
        if "/" in piece:
            rng, step_s = piece.split("/")
            step = int(step_s)
            if rng == "*":
                start, end = lo, hi
            elif "-" in rng:
                a, b = rng.split("-")
                start, end = int(a), int(b)
            else:
                start, end = int(rng), hi
            out.update(range(start, end + 1, step))
        elif "-" in piece:
            a, b = piece.split("-")
            out.update(range(int(a), int(b) + 1))
        else:
            out.add(int(piece))
    return sorted(i for i in out if lo <= i <= hi)


def cron_next(expression: str, now: datetime) -> datetime:
    mins, hours, doms, months, dows = expression.split()
    m_set = _parse_cron_field(mins, 0, 59)
    h_set = _parse_cron_field(hours, 0, 23)
    dom_set = _parse_cron_field(doms, 1, 31)
    mo_set = _parse_cron_field(months, 1, 12)
    dow_set = _parse_cron_field(dows, 0, 6)
    t = now.replace(second=0, microsecond=0) + timedelta(minutes=1)
    for _ in range(60 * 24 * 366):
        if (
            t.minute in m_set
            and t.hour in h_set
            and t.day in dom_set
            and t.month in mo_set
            and (t.weekday() + 1) % 7 in dow_set
        ):
            return t
        t += timedelta(minutes=1)
    raise RuntimeError("no cron match within a year")


# ---------------------------------------------------------------------------
# 9. Dispatch mock
# ---------------------------------------------------------------------------

@dataclass
class Dispatch:
    queue: List[Dict[str, Any]] = field(default_factory=list)

    def send(self, device: str, task: str) -> str:
        msg_id = str(uuid.uuid4())
        self.queue.append({"id": msg_id, "device": device, "task": task, "status": "queued"})
        return msg_id

    def execute_next(self) -> Optional[Dict[str, Any]]:
        for msg in self.queue:
            if msg["status"] == "queued":
                msg["status"] = "done"
                return msg
        return None


# ---------------------------------------------------------------------------
# 10. Connector mocks
# ---------------------------------------------------------------------------

class MockConnector:
    def __init__(self, name: str):
        self.name = name
        self.items: List[Dict[str, Any]] = []

    def upsert(self, item: Dict[str, Any]) -> None:
        for i, existing in enumerate(self.items):
            if existing.get("id") == item.get("id"):
                self.items[i] = item
                return
        self.items.append(item)

    def list(self) -> List[Dict[str, Any]]:
        return list(self.items)

    def search(self, term: str) -> List[Dict[str, Any]]:
        tl = term.lower()
        return [i for i in self.items if tl in json.dumps(i).lower()]


CONNECTOR_NAMES = [
    "google_drive", "gmail", "docusign", "zoom", "slack", "jira",
    "asana", "notion", "github", "linear", "hubspot", "fellow",
]


# ---------------------------------------------------------------------------
# 11. Plugins / skill-creator mock
# ---------------------------------------------------------------------------

@dataclass
class Skill:
    name: str
    description: str
    triggers: List[str]

    def matches(self, prompt: str) -> bool:
        p = prompt.lower()
        return any(t.lower() in p for t in self.triggers)


def create_skill(name: str, description: str, triggers: Iterable[str]) -> Skill:
    return Skill(name=name, description=description, triggers=list(triggers))


# ---------------------------------------------------------------------------
# 12. Code execution sandbox
# ---------------------------------------------------------------------------

def exec_python(code: str, *, timeout: float = 5.0) -> Tuple[int, str, str]:
    with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False) as f:
        f.write(code)
        fname = f.name
    try:
        proc = subprocess.run(
            [sys.executable, fname],
            capture_output=True, text=True, timeout=timeout,
        )
        return proc.returncode, proc.stdout, proc.stderr
    finally:
        os.unlink(fname)


def exec_node(code: str, *, timeout: float = 5.0) -> Tuple[int, str, str]:
    with tempfile.NamedTemporaryFile("w", suffix=".js", delete=False) as f:
        f.write(code)
        fname = f.name
    try:
        proc = subprocess.run(
            ["node", fname],
            capture_output=True, text=True, timeout=timeout,
        )
        return proc.returncode, proc.stdout, proc.stderr
    finally:
        os.unlink(fname)


# ---------------------------------------------------------------------------
# 13. Sub-agents / task decomposition
# ---------------------------------------------------------------------------

def decompose_task(task: str) -> List[str]:
    """Rule-based decomposition; deterministic for tests."""
    t = task.lower()
    steps: List[str] = []
    if "report" in t:
        steps += ["gather sources", "outline", "draft", "review", "export"]
    if "analy" in t:
        steps += ["load data", "clean", "explore", "model", "summarize"]
    if "email" in t:
        steps += ["read inbox", "triage", "draft replies", "review"]
    if not steps:
        steps = ["understand", "plan", "execute", "verify"]
    return steps


class TodoList:
    def __init__(self) -> None:
        self.items: List[Dict[str, Any]] = []

    def add(self, content: str) -> int:
        idx = len(self.items)
        self.items.append({"id": idx, "content": content, "status": "pending"})
        return idx

    def start(self, idx: int) -> None:
        self.items[idx]["status"] = "in_progress"

    def done(self, idx: int) -> None:
        self.items[idx]["status"] = "completed"

    def progress(self) -> float:
        if not self.items:
            return 0.0
        done = sum(1 for i in self.items if i["status"] == "completed")
        return done / len(self.items)


# ---------------------------------------------------------------------------
# 14. Project workspaces
# ---------------------------------------------------------------------------

@dataclass
class Workspace:
    name: str
    root: pathlib.Path
    memory: Dict[str, Any] = field(default_factory=dict)

    def save(self) -> pathlib.Path:
        state_path = self.root / ".workspace.json"
        state_path.parent.mkdir(parents=True, exist_ok=True)
        state_path.write_text(json.dumps({"name": self.name, "memory": self.memory}))
        return state_path

    @classmethod
    def load(cls, root: pathlib.Path) -> "Workspace":
        state = json.loads((root / ".workspace.json").read_text())
        return cls(name=state["name"], root=root, memory=state.get("memory", {}))


# ---------------------------------------------------------------------------
# 15. Security / governance
# ---------------------------------------------------------------------------

def path_within(root: pathlib.Path, target: pathlib.Path) -> bool:
    try:
        target.resolve().relative_to(root.resolve())
        return True
    except ValueError:
        return False


def egress_allowed(url: str, allowlist: Sequence[str]) -> bool:
    host = urllib.parse.urlparse(url).hostname or ""
    return any(host == h or host.endswith("." + h) for h in allowlist)


# ---------------------------------------------------------------------------
# 16. Enterprise
# ---------------------------------------------------------------------------

class RBAC:
    def __init__(self) -> None:
        self.roles: Dict[str, set[str]] = {}
        self.user_roles: Dict[str, set[str]] = {}

    def grant(self, role: str, capability: str) -> None:
        self.roles.setdefault(role, set()).add(capability)

    def assign(self, user: str, role: str) -> None:
        self.user_roles.setdefault(user, set()).add(role)

    def can(self, user: str, capability: str) -> bool:
        for role in self.user_roles.get(user, set()):
            if capability in self.roles.get(role, set()):
                return True
        return False


class Budget:
    def __init__(self, limit_cents: int) -> None:
        self.limit = limit_cents
        self.spent = 0

    def charge(self, amount_cents: int) -> bool:
        if self.spent + amount_cents > self.limit:
            return False
        self.spent += amount_cents
        return True


def telemetry_aggregate(events: Sequence[Dict[str, Any]]) -> Dict[str, int]:
    agg: Dict[str, int] = {}
    for e in events:
        k = e.get("event", "unknown")
        agg[k] = agg.get(k, 0) + 1
    return agg


# ---------------------------------------------------------------------------
# 17. Domain templates
# ---------------------------------------------------------------------------

DOMAIN_TEMPLATES: Dict[str, Dict[str, List[str]]] = {
    "legal": {
        "contract_review": ["parties", "term", "termination", "indemnity", "governing law"],
        "nda_triage": ["confidential info", "duration", "exclusions", "remedies"],
    },
    "finance": {
        "month_close": ["journal entries", "reconciliation", "variance analysis", "sign off"],
        "budget": ["forecast", "actual", "variance", "commentary"],
    },
    "marketing": {
        "brand_voice": ["tone", "vocabulary", "dos", "donts"],
    },
    "ops": {
        "daily_brief": ["agenda", "priorities", "risks", "owners"],
    },
    "hr": {
        "perf_review": ["goals", "achievements", "areas", "next"],
    },
    "research": {
        "interview_synthesis": ["themes", "quotes", "pain points", "actions"],
    },
}


def render_template(domain: str, name: str, values: Dict[str, str]) -> str:
    sections = DOMAIN_TEMPLATES[domain][name]
    out = [f"# {domain.title()} — {name.replace('_', ' ').title()}"]
    for s in sections:
        out.append(f"## {s.title()}")
        out.append(values.get(s, "(tbd)"))
    return "\n".join(out)
