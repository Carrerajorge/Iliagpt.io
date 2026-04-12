"""
CAP-02: Generacion de Archivos PowerPoint (.pptx)
==================================================

Validates PPTX generation: from scratch, with speaker notes, watermarks,
document-to-slides conversion, editability after generation.

Sub-capabilities under test:
  2.1  Presentaciones desde cero o desde notas/transcripciones
  2.2  Layouts con speaker notes
  2.3  Conversion de documentos a slides
  2.4  GIFs/imagenes (placeholder validation)
  2.5  Watermark en lote
  2.6  Editable despues de generacion
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from cowork_lib3 import (
    create_pptx_with_notes,
    create_pptx_with_watermark,
    create_pptx_from_doc_outline,
    create_pptx_multi_layout,
)
from cowork_lib import generate_document, DocSpec, read_document_text


@pytest.fixture
def out_dir(tmp_path):
    d = tmp_path / "cap02_pptx"; d.mkdir(); return d


@pytest.mark.file_generation
class TestPptxFromScratch:
    """2.1 — Presentation creation from prompts and transcripts."""

    @pytest.mark.parametrize("n_slides", [1, 3, 5, 10, 20])
    def test_slide_count(self, out_dir, n_slides):
        slides = [{"title": f"S{i}", "body": f"C{i}", "notes": f"N{i}"} for i in range(n_slides)]
        path = out_dir / f"scratch_{n_slides}.pptx"
        create_pptx_with_notes(path, slides)
        assert len(Presentation(path).slides) == n_slides

    @pytest.mark.parametrize("i", range(20))
    def test_title_preserved(self, out_dir, i):
        slides = [{"title": f"Title_{i}", "body": f"Body_{i}", "notes": ""}]
        path = out_dir / f"title_{i}.pptx"
        create_pptx_with_notes(path, slides)
        assert Presentation(path).slides[0].shapes.title.text == f"Title_{i}"

    @pytest.mark.parametrize("i", range(15))
    def test_transcript_to_slides(self, out_dir, i):
        sections = [f"Meeting Topic {i}", "Item 1: Alpha", "Item 2: Budget", "Item 3: Next Steps", "Conclusion"]
        slides = [{"title": s, "body": f"Details: {s}", "notes": f"Note: {s}"} for s in sections]
        path = out_dir / f"transcript_{i}.pptx"
        create_pptx_with_notes(path, slides)
        prs = Presentation(path)
        assert len(prs.slides) == 5
        assert "Meeting Topic" in prs.slides[0].shapes.title.text


@pytest.mark.file_generation
class TestPptxSpeakerNotes:
    """2.2 — Speaker notes attached to every slide."""

    @pytest.mark.parametrize("i", range(20))
    def test_notes_content(self, out_dir, i):
        slides = [
            {"title": f"Intro {i}", "body": "Welcome", "notes": f"Greet audience {i}"},
            {"title": "Main", "body": "Details", "notes": f"Key metric: {i*100}"},
        ]
        path = out_dir / f"notes_{i}.pptx"
        create_pptx_with_notes(path, slides)
        prs = Presentation(path)
        for j, slide in enumerate(prs.slides):
            assert slide.notes_slide.notes_text_frame.text == slides[j]["notes"]

    @pytest.mark.parametrize("n", range(1, 11))
    def test_all_slides_have_notes(self, out_dir, n):
        slides = [{"title": f"S{i}", "body": "", "notes": f"Note {i}"} for i in range(n)]
        path = out_dir / f"all_notes_{n}.pptx"
        create_pptx_with_notes(path, slides)
        for slide in Presentation(path).slides:
            assert slide.notes_slide.notes_text_frame.text != ""

    @pytest.mark.parametrize("i", range(10))
    def test_multi_layout_types(self, out_dir, i):
        content = [
            {"layout": 0, "title": f"Title {i}", "notes": "Opening"},
            {"layout": 1, "title": f"Content {i}", "notes": "Body"},
            {"layout": 5, "title": f"Blank {i}", "notes": "Visual"},
        ]
        path = out_dir / f"layouts_{i}.pptx"
        create_pptx_multi_layout(path, content)
        assert len(Presentation(path).slides) == 3


@pytest.mark.file_generation
class TestPptxDocConversion:
    """2.3 — Document paragraphs converted to slide deck."""

    @pytest.mark.parametrize("n_paragraphs", [3, 5, 8, 12])
    def test_paragraph_to_slide_count(self, out_dir, n_paragraphs):
        paragraphs = [f"Heading {i}: Topic {i}" for i in range(n_paragraphs)]
        path = out_dir / f"doc2slides_{n_paragraphs}.pptx"
        create_pptx_from_doc_outline(paragraphs, path)
        assert len(Presentation(path).slides) == n_paragraphs

    @pytest.mark.parametrize("i", range(10))
    def test_title_from_first_paragraph(self, out_dir, i):
        paragraphs = [f"Report Title {i}", "Intro", "Body", "Conclusion"]
        path = out_dir / f"outline_{i}.pptx"
        create_pptx_from_doc_outline(paragraphs, path)
        assert Presentation(path).slides[0].shapes.title.text == f"Report Title {i}"

    @pytest.mark.parametrize("i", range(10))
    def test_full_docx_to_pptx_pipeline(self, out_dir, i):
        spec = DocSpec(kind="docx", title=f"Doc {i}", sections=[f"Section {j}" for j in range(4)])
        docx_path = generate_document(spec, out_dir)
        text = read_document_text(docx_path)
        paragraphs = [p for p in text.split("\n") if p.strip()]
        pptx_path = out_dir / f"from_docx_{i}.pptx"
        create_pptx_from_doc_outline(paragraphs, pptx_path)
        assert len(Presentation(pptx_path).slides) >= 2


@pytest.mark.file_generation
class TestPptxImagePlaceholders:
    """2.4 — GIF/image placeholder slots on slides."""

    @pytest.mark.parametrize("i", range(15))
    def test_placeholder_slot_exists(self, out_dir, i):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Image Slide {i}"
        assert len(list(slide.placeholders)) >= 1
        path = out_dir / f"ph_{i}.pptx"
        prs.save(path)
        assert path.stat().st_size > 0

    @pytest.mark.parametrize("i", range(10))
    def test_gif_marker_textbox(self, out_dir, i):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
        txBox.text_frame.paragraphs[0].text = f"[GIF_PLACEHOLDER_{i}]"
        path = out_dir / f"gif_{i}.pptx"
        prs.save(path)
        found = any("GIF_PLACEHOLDER" in s.text for s in Presentation(path).slides[0].shapes if hasattr(s, "text"))
        assert found


WATERMARKS = ["CONFIDENTIAL", "DRAFT", "INTERNAL", "DO NOT DISTRIBUTE", "PRIVILEGED"]


@pytest.mark.file_generation
class TestPptxWatermark:
    """2.5 — Batch watermark on every slide."""

    @pytest.mark.parametrize("wm", WATERMARKS)
    def test_watermark_text_present(self, out_dir, wm):
        path = out_dir / f"wm_{wm.replace(' ','_')}.pptx"
        create_pptx_with_watermark(path, ["S1", "S2", "S3"], watermark=wm)
        for slide in Presentation(path).slides:
            assert any(wm in s.text for s in slide.shapes if hasattr(s, "text"))

    @pytest.mark.parametrize("n_slides", [1, 5, 10, 20])
    def test_watermark_on_every_slide(self, out_dir, n_slides):
        path = out_dir / f"wm_all_{n_slides}.pptx"
        create_pptx_with_watermark(path, [f"C{i}" for i in range(n_slides)], watermark="DRAFT")
        prs = Presentation(path)
        assert len(prs.slides) == n_slides
        for slide in prs.slides:
            assert any("DRAFT" in s.text for s in slide.shapes if hasattr(s, "text"))

    @pytest.mark.parametrize("i", range(10))
    def test_watermark_count_per_slide(self, out_dir, i):
        path = out_dir / f"wm_batch_{i}.pptx"
        create_pptx_with_watermark(path, [f"S{j}" for j in range(5)], watermark="CONFIDENTIAL")
        wm_count = sum(
            1 for slide in Presentation(path).slides
            for s in slide.shapes if hasattr(s, "text") and "CONFIDENTIAL" in s.text
        )
        assert wm_count == 5


@pytest.mark.file_generation
class TestPptxEditability:
    """2.6 — Generated PPTX files must be re-openable and editable."""

    @pytest.mark.parametrize("i", range(15))
    def test_modify_after_creation(self, out_dir, i):
        slides = [{"title": f"Original {i}", "body": "C", "notes": "N"}]
        path = out_dir / f"editable_{i}.pptx"
        create_pptx_with_notes(path, slides)
        prs = Presentation(path)
        prs.slides[0].shapes.title.text = f"Modified {i}"
        prs.slides.add_slide(prs.slide_layouts[1])
        mod_path = out_dir / f"modified_{i}.pptx"
        prs.save(mod_path)
        prs2 = Presentation(mod_path)
        assert len(prs2.slides) == 2
        assert prs2.slides[0].shapes.title.text == f"Modified {i}"

    @pytest.mark.parametrize("i", range(10))
    def test_modify_notes(self, out_dir, i):
        path = out_dir / f"mn_{i}.pptx"
        create_pptx_with_notes(path, [{"title": "T", "body": "B", "notes": f"Old {i}"}])
        prs = Presentation(path)
        prs.slides[0].notes_slide.notes_text_frame.text = f"New {i}"
        p2 = out_dir / f"mn_saved_{i}.pptx"
        prs.save(p2)
        assert Presentation(p2).slides[0].notes_slide.notes_text_frame.text == f"New {i}"
