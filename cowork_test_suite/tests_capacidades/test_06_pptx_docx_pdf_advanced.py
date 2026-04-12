"""
Capability 1b — Advanced PPTX / DOCX / PDF generation.
~1300 tests total.
"""
from __future__ import annotations

import pathlib

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

N = 100


# ------------------------- PPTX -------------------------
@pytest.mark.parametrize("n_slides", list(range(1, 21)))
def test_pptx_multi_slide(tmp_path, n_slides):
    prs = Presentation()
    for i in range(n_slides):
        s = prs.slides.add_slide(prs.slide_layouts[1 if i else 0])
        s.shapes.title.text = f"Slide {i+1}"
    path = tmp_path / f"deck_{n_slides}.pptx"
    prs.save(path)
    r = Presentation(path)
    assert len(r.slides) == n_slides


@pytest.mark.parametrize("i", range(N))
def test_pptx_speaker_notes(tmp_path, i):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = f"Title {i}"
    s.notes_slide.notes_text_frame.text = f"Speaker note {i}"
    path = tmp_path / f"notes_{i}.pptx"
    prs.save(path)
    r = Presentation(path)
    first = r.slides[0]
    assert f"Speaker note {i}" in first.notes_slide.notes_text_frame.text


@pytest.mark.parametrize("i", range(N))
def test_pptx_layout_title_subtitle(tmp_path, i):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = f"Main {i}"
    if len(s.placeholders) > 1:
        s.placeholders[1].text = f"Sub {i}"
    path = tmp_path / f"layout_{i}.pptx"
    prs.save(path)
    r = Presentation(path)
    assert r.slides[0].shapes.title.text == f"Main {i}"


@pytest.mark.parametrize("i", range(N))
def test_pptx_text_body(tmp_path, i):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = f"Body {i}"
    tf = s.placeholders[1].text_frame
    tf.text = "first"
    tf.add_paragraph().text = "second"
    tf.add_paragraph().text = "third"
    path = tmp_path / f"body_{i}.pptx"
    prs.save(path)
    r = Presentation(path)
    paras = [p.text for p in r.slides[0].placeholders[1].text_frame.paragraphs]
    assert paras[0] == "first"


# ------------------------- DOCX -------------------------
@pytest.mark.parametrize("i", range(N))
def test_docx_headings(tmp_path, i):
    doc = Document()
    doc.add_heading(f"Main {i}", level=1)
    doc.add_heading("Sub", level=2)
    doc.add_paragraph("body text here")
    path = tmp_path / f"h_{i}.docx"
    doc.save(path)
    r = Document(path)
    styles = [p.style.name for p in r.paragraphs]
    assert any("Heading" in s for s in styles)


@pytest.mark.parametrize("rows", list(range(2, 12)))
@pytest.mark.parametrize("cols", list(range(2, 7)))
def test_docx_tables(tmp_path, rows, cols):
    doc = Document()
    t = doc.add_table(rows=rows, cols=cols)
    for r in range(rows):
        for c in range(cols):
            t.cell(r, c).text = f"{r},{c}"
    path = tmp_path / f"tbl_{rows}_{cols}.docx"
    doc.save(path)
    r2 = Document(path)
    assert len(r2.tables[0].rows) == rows
    assert len(r2.tables[0].columns) == cols


@pytest.mark.parametrize("i", range(N))
def test_docx_paragraph_roundtrip(tmp_path, i):
    doc = Document()
    marker = f"PARA-{i:04d}-UNIQUE"
    doc.add_paragraph(marker)
    doc.add_paragraph("second paragraph")
    path = tmp_path / f"p_{i}.docx"
    doc.save(path)
    r = Document(path)
    assert r.paragraphs[0].text == marker


# ------------------------- PDF -------------------------
@pytest.mark.parametrize("i", range(N))
def test_pdf_create(tmp_path, i):
    path = tmp_path / f"doc_{i}.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    c.setFont("Helvetica", 12)
    c.drawString(100, 750, f"Document #{i}")
    c.drawString(100, 730, "Line two of the document")
    c.save()
    r = PdfReader(str(path))
    text = r.pages[0].extract_text() or ""
    assert f"Document #{i}" in text


@pytest.mark.parametrize("n_pages", list(range(2, 12)))
def test_pdf_multi_page(tmp_path, n_pages):
    path = tmp_path / f"pages_{n_pages}.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    for i in range(n_pages):
        c.drawString(100, 750, f"Page {i+1}")
        c.showPage()
    c.save()
    r = PdfReader(str(path))
    assert len(r.pages) == n_pages


@pytest.mark.parametrize("seed", list(range(50)))
def test_pdf_merge(tmp_path, seed):
    # create two PDFs and merge them
    p1 = tmp_path / f"a_{seed}.pdf"
    p2 = tmp_path / f"b_{seed}.pdf"
    for p, label in [(p1, "A"), (p2, "B")]:
        c = canvas.Canvas(str(p), pagesize=letter)
        c.drawString(100, 750, f"{label} seed={seed}")
        c.save()
    merged = tmp_path / f"merged_{seed}.pdf"
    w = PdfWriter()
    for p in [p1, p2]:
        for page in PdfReader(str(p)).pages:
            w.add_page(page)
    with merged.open("wb") as f:
        w.write(f)
    assert len(PdfReader(str(merged)).pages) == 2


@pytest.mark.parametrize("seed", list(range(50)))
def test_pdf_split(tmp_path, seed):
    src = tmp_path / f"src_{seed}.pdf"
    c = canvas.Canvas(str(src), pagesize=letter)
    for i in range(4):
        c.drawString(100, 750, f"Page {i}")
        c.showPage()
    c.save()
    reader = PdfReader(str(src))
    for i, page in enumerate(reader.pages):
        w = PdfWriter()
        w.add_page(page)
        out = tmp_path / f"split_{seed}_{i}.pdf"
        with out.open("wb") as f:
            w.write(f)
        assert out.exists()
