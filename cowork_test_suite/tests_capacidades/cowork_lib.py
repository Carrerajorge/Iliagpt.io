"""
cowork_lib
==========

A minimal "Cowork-like" library implementing the four capability pillars the
user asked us to validate:

1. Document generation  (docx / xlsx / pptx / pdf)
2. Information search   (full-text + hybrid scoring)
3. Information validation (real-world fact patterns: dates, URLs, emails,
   ISBNs, credit cards Luhn, UUIDs, JSON structure)
4. Code generation      (template-driven code emission + syntax check)

Everything here is deterministic and self-contained so that thousands of
parametrized tests can exercise it without network, database or LLM access.
"""
from __future__ import annotations

import ast
import datetime as dt
import hashlib
import json
import math
import re
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable, Sequence

# ---------------------------------------------------------------------------
# 1. Document generation
# ---------------------------------------------------------------------------

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pypdf import PdfReader


@dataclass
class DocSpec:
    """Declarative description of a document to generate."""
    kind: str  # "docx" | "xlsx" | "pptx" | "pdf"
    title: str
    sections: list[str] = field(default_factory=list)
    data: list[list[Any]] | None = None


def generate_document(spec: DocSpec, out_dir: Path) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    safe_title = re.sub(r"[^a-zA-Z0-9_-]+", "_", spec.title)[:60]
    path = out_dir / f"{safe_title}.{spec.kind}"

    if spec.kind == "docx":
        doc = Document()
        doc.add_heading(spec.title, level=1)
        for s in spec.sections:
            doc.add_paragraph(s)
        doc.save(path)

    elif spec.kind == "xlsx":
        wb = Workbook()
        ws = wb.active
        ws.title = spec.title[:30] or "Sheet"
        rows = spec.data or [["A", "B", "C"], [1, 2, 3]]
        for r in rows:
            ws.append(r)
        wb.save(path)

    elif spec.kind == "pptx":
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = spec.title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = spec.sections[0] if spec.sections else ""
        for sec in spec.sections[1:]:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = sec[:60]
        prs.save(path)

    elif spec.kind == "pdf":
        c = canvas.Canvas(str(path), pagesize=letter)
        c.setFont("Helvetica-Bold", 18)
        c.drawString(72, 720, spec.title)
        c.setFont("Helvetica", 11)
        y = 690
        for s in spec.sections:
            c.drawString(72, y, s[:90])
            y -= 16
            if y < 80:
                c.showPage()
                y = 720
        c.save()
    else:
        raise ValueError(f"unknown kind {spec.kind}")

    return path


def read_document_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        d = Document(path)
        return "\n".join(p.text for p in d.paragraphs)
    if suffix == ".xlsx":
        wb = load_workbook(path)
        ws = wb.active
        return "\n".join(
            " ".join("" if c is None else str(c) for c in row)
            for row in ws.iter_rows(values_only=True)
        )
    if suffix == ".pptx":
        prs = Presentation(path)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        out.append(para.text)
        return "\n".join(out)
    if suffix == ".pdf":
        r = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in r.pages)
    raise ValueError(suffix)


# ---------------------------------------------------------------------------
# 2. Information search
# ---------------------------------------------------------------------------

_WORD_RE = re.compile(r"\w+", re.UNICODE)


def tokenize(text: str) -> list[str]:
    return [t.lower() for t in _WORD_RE.findall(text)]


class SearchIndex:
    """Small TF-IDF style index with hybrid BM25-ish scoring."""

    def __init__(self, docs: Sequence[str]):
        self.docs = list(docs)
        self.tokens = [tokenize(d) for d in self.docs]
        self.df: dict[str, int] = {}
        for toks in self.tokens:
            for w in set(toks):
                self.df[w] = self.df.get(w, 0) + 1
        self.N = len(self.docs)
        self.avgdl = (sum(len(t) for t in self.tokens) / self.N) if self.N else 0

    def idf(self, term: str) -> float:
        n = self.df.get(term, 0)
        return math.log(1 + (self.N - n + 0.5) / (n + 0.5))

    def score(self, query: str, doc_index: int, k1: float = 1.5, b: float = 0.75) -> float:
        q = tokenize(query)
        toks = self.tokens[doc_index]
        if not toks:
            return 0.0
        dl = len(toks)
        freqs: dict[str, int] = {}
        for t in toks:
            freqs[t] = freqs.get(t, 0) + 1
        s = 0.0
        for term in q:
            f = freqs.get(term, 0)
            if f == 0:
                continue
            num = f * (k1 + 1)
            den = f + k1 * (1 - b + b * dl / (self.avgdl or 1))
            s += self.idf(term) * num / den
        return s

    def search(self, query: str, top_k: int = 5) -> list[tuple[int, float]]:
        scored = [(i, self.score(query, i)) for i in range(self.N)]
        scored = [p for p in scored if p[1] > 0]
        scored.sort(key=lambda x: x[1], reverse=True)
        return scored[:top_k]


# ---------------------------------------------------------------------------
# 3. Real-world information validation
# ---------------------------------------------------------------------------

EMAIL_RE = re.compile(
    r"^[A-Z0-9][A-Z0-9._%+-]*@[A-Z0-9]([A-Z0-9-]*[A-Z0-9])?(\.[A-Z0-9]([A-Z0-9-]*[A-Z0-9])?)*\.[A-Z]{2,}$",
    re.IGNORECASE,
)
URL_RE = re.compile(
    r"^https?://([A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+[A-Z]{2,}(/[^\s]*)?$",
    re.IGNORECASE,
)
UUID_RE = re.compile(r"^[0-9a-f]{8}-[0-9a-f]{4}-[1-5][0-9a-f]{3}-[89ab][0-9a-f]{3}-[0-9a-f]{12}$", re.IGNORECASE)


def is_email(s: str) -> bool:
    return bool(EMAIL_RE.match(s))


def is_url(s: str) -> bool:
    return bool(URL_RE.match(s))


def is_uuid(s: str) -> bool:
    return bool(UUID_RE.match(s))


def is_iso_date(s: str) -> bool:
    try:
        dt.date.fromisoformat(s)
        return True
    except Exception:
        return False


def luhn_check(card: str) -> bool:
    digits = [int(c) for c in card if c.isdigit()]
    if len(digits) < 12:
        return False
    total = 0
    parity = len(digits) % 2
    for i, d in enumerate(digits):
        if i % 2 == parity:
            d *= 2
            if d > 9:
                d -= 9
        total += d
    return total % 10 == 0


def is_isbn10(s: str) -> bool:
    s = s.replace("-", "").replace(" ", "")
    if len(s) != 10:
        return False
    total = 0
    for i, c in enumerate(s):
        if c == "X" and i == 9:
            total += 10 * 1
            continue
        if not c.isdigit():
            return False
        total += int(c) * (10 - i)
    return total % 11 == 0


def is_isbn13(s: str) -> bool:
    s = s.replace("-", "").replace(" ", "")
    if len(s) != 13 or not s.isdigit():
        return False
    total = sum(int(c) * (1 if i % 2 == 0 else 3) for i, c in enumerate(s))
    return total % 10 == 0


def validate_json(s: str) -> bool:
    try:
        json.loads(s)
        return True
    except Exception:
        return False


def sha256_hex(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8")).hexdigest()


# ---------------------------------------------------------------------------
# 4. Code generation
# ---------------------------------------------------------------------------

PY_FUNC_TEMPLATE = '''\
def {name}({args}):
    """{doc}"""
    {body}
'''


def generate_python_function(name: str, args: Iterable[str], body: str, doc: str = "") -> str:
    code = PY_FUNC_TEMPLATE.format(
        name=name,
        args=", ".join(args),
        doc=doc or f"Auto-generated {name}",
        body=body,
    )
    return code


def python_syntax_ok(code: str) -> bool:
    try:
        ast.parse(code)
        return True
    except SyntaxError:
        return False


def run_python(code: str, ns: dict | None = None) -> dict:
    ns = ns if ns is not None else {}
    exec(compile(code, "<cowork>", "exec"), ns)
    return ns


# Deterministic corpus used by the search tests ------------------------------

CORPUS = [
    "Python is a programming language that lets you work quickly.",
    "The report on quarterly sales shows strong growth in Europe.",
    "Machine learning models require clean training data and evaluation.",
    "Climate change is causing rising sea levels and extreme weather.",
    "The conference will be held in Berlin in September 2026.",
    "Electric vehicles dominated the auto show in Tokyo last week.",
    "Quantum computing promises exponential speedups for certain problems.",
    "The recipe calls for two cups of flour, sugar, and vanilla extract.",
    "Renewable energy accounts for more than 30% of global electricity.",
    "The novel won the Booker Prize for its innovative structure.",
    "SpaceX launched a new batch of Starlink satellites this morning.",
    "The vaccine trial reported 95% efficacy in phase three results.",
    "Blockchain technology underpins cryptocurrencies and smart contracts.",
    "The hiking trail passes through alpine meadows and pine forests.",
    "Artificial intelligence is transforming the healthcare industry.",
    "The museum exhibit features ancient Egyptian artifacts and mummies.",
    "Coffee cultivation depends heavily on equatorial climates.",
    "The author published a new book on economic policy and reform.",
    "Tropical cyclones form over warm ocean waters near the equator.",
    "The startup raised twenty million dollars in Series B funding.",
    "A new species of frog was discovered in the Amazon rainforest.",
    "The chess grandmaster defeated the reigning world champion.",
    "Solar panels on residential homes have become more affordable.",
    "Archaeologists uncovered a Roman mosaic near the ruins of Pompeii.",
    "The orchestra performed a symphony by Gustav Mahler last night.",
]
