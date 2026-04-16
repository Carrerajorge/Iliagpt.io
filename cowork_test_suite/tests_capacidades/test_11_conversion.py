"""
Capability 5 — Format conversion: md->html, csv->xlsx model, docx->pptx, pptx->md.
~400 tests.
"""
from __future__ import annotations

import pathlib

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from cowork_lib2 import (
    markdown_to_html, csv_to_xlsx_model, docx_to_pptx_outline,
    pptx_to_markdown_outline, write_csv,
)


# ---- markdown -> html ----
@pytest.mark.parametrize("i", range(100))
def test_md_to_html(i):
    md = f"# Title {i}\n\nSome **bold** text and a list:\n\n- one\n- two\n"
    h = markdown_to_html(md)
    assert f"Title {i}" in h
    assert "<ul" in h.lower()


# ---- csv -> xlsx model ----
@pytest.mark.parametrize("i", range(100))
def test_csv_to_xlsx(tmp_path, i):
    rows = [["item", "price"]] + [[f"p{j}", j * 10] for j in range(1, 6)]
    src = write_csv(tmp_path / f"src_{i}.csv", rows)
    dst = tmp_path / f"out_{i}.xlsx"
    csv_to_xlsx_model(src, dst)
    wb = load_workbook(dst)
    ws = wb.active
    assert ws.cell(row=1, column=1).value == "item"
    # total row present
    total_row = ws.max_row
    assert ws.cell(row=total_row, column=1).value == "TOTAL"
    assert "SUM" in ws.cell(row=total_row, column=2).value


# ---- docx -> pptx outline ----
@pytest.mark.parametrize("i", range(100))
def test_docx_to_pptx(tmp_path, i):
    doc = Document()
    doc.add_paragraph(f"Deck {i}")
    for j in range(3):
        doc.add_paragraph(f"Point {j}")
    d = tmp_path / f"in_{i}.docx"
    doc.save(d)
    out = tmp_path / f"out_{i}.pptx"
    docx_to_pptx_outline(d, out)
    prs = Presentation(out)
    assert len(prs.slides) == 4  # title + 3 points


# ---- pptx -> markdown outline ----
@pytest.mark.parametrize("i", range(100))
def test_pptx_to_md(tmp_path, i):
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = f"Deck {i}"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Section"
    p = tmp_path / f"deck_{i}.pptx"
    prs.save(p)
    md = pptx_to_markdown_outline(p)
    assert f"Deck {i}" in md
    assert "Slide 1" in md and "Slide 2" in md
